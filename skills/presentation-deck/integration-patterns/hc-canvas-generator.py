"""
Health Canada Canvas Generator
Programmatically create branded HC/PHAC canvas slides (Opportunity Canvas, Business Model Canvas, etc.)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass
import json

# HC/PHAC Brand Colors
HC_COLORS = {
    "primary": RGBColor(0, 61, 104),        # Dark Blue
    "secondary": RGBColor(237, 246, 255),  # Light Blue
    "tertiary": RGBColor(229, 255, 239),   # Light Green
    "neutral": RGBColor(247, 249, 251),    # Light Gray
    "border": RGBColor(148, 163, 184),     # Slate
    "text_dark": RGBColor(17, 24, 39),     # Dark Gray
    "text_white": RGBColor(255, 255, 255), # White
}


@dataclass
class Card:
    """Represents a single canvas card"""
    title: str
    lines: List[str]
    x: float  # in inches
    y: float  # in inches
    width: float  # in inches
    height: float  # in inches
    fill_color: RGBColor
    title_color: RGBColor = HC_COLORS["primary"]


class CanvasBuilder:
    """Builder for HC-branded canvas slides"""

    def __init__(self, title: str, width: float = 17, height: float = 11):
        """
        Initialize canvas

        Args:
            title: Canvas title (appears in header)
            width: Slide width in inches
            height: Slide height in inches
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.title = title
        self._add_header()

    def _add_header(self):
        """Add HC-branded header"""
        header = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(17), Inches(1.0)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = HC_COLORS["primary"]
        header.line.fill.background()

        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = HC_COLORS["text_white"]

    def add_card(
        self,
        title: str,
        lines: List[str],
        x: float,
        y: float,
        width: float = 5.4,
        height: float = 1.6,
        fill_color: RGBColor = HC_COLORS["neutral"],
        title_color: RGBColor = HC_COLORS["primary"]
    ):
        """
        Add a card to the canvas

        Args:
            title: Card title (bold)
            lines: List of content lines
            x, y: Position in inches
            width, height: Card dimensions in inches
            fill_color: Background color (RGB)
            title_color: Title text color (RGB)
        """
        card = Card(
            title=title,
            lines=lines,
            x=Inches(x),
            y=Inches(y),
            width=Inches(width),
            height=Inches(height),
            fill_color=fill_color,
            title_color=title_color
        )
        self._render_card(card)

    def _render_card(self, card: Card):
        """Render a card on the slide"""
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            card.x, card.y,
            card.width, card.height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = card.fill_color
        shape.line.color.rgb = HC_COLORS["border"]
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.18)
        tf.margin_bottom = Inches(0.18)

        # Title
        p = tf.paragraphs[0]
        p.text = card.title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = card.title_color

        # Content lines
        for line in card.lines:
            q = tf.add_paragraph()
            q.text = line
            q.font.size = Pt(14)
            q.font.color.rgb = HC_COLORS["text_dark"]

    def save(self, path: str = "canvas.pptx"):
        """Save presentation to file"""
        self.prs.save(path)
        return path


class OpportunityCanvas(CanvasBuilder):
    """Lean Canvas / Opportunity Canvas builder"""

    def __init__(self, title: str = "Opportunity Canvas"):
        super().__init__(title)

    def build(self, data: Dict) -> str:
        """
        Build opportunity canvas from data dict

        Args:
            data: Dictionary with keys: target_users, problems, alternatives,
                  contexts, solution, value_prop, adoption, metrics,
                  business_value, costs, risks, questions

        Returns:
            Path to saved PPTX
        """
        # Column positions
        x1, x2, x3 = 0.3, 5.9, 11.5
        y = 1.3

        # Left column (problems/context - neutral gray)
        self.add_card("Target users and customers", data.get("target_users", []),
                     x1, y, fill_color=HC_COLORS["neutral"])
        self.add_card("User problems and needs", data.get("problems", []),
                     x1, y + 1.75, fill_color=HC_COLORS["neutral"])
        self.add_card("Current alternatives", data.get("alternatives", []),
                     x1, y + 3.5, fill_color=HC_COLORS["neutral"])
        self.add_card("Usage contexts", data.get("contexts", []),
                     x1, y + 5.25, Inches(5.4), Inches(2.25),
                     fill_color=HC_COLORS["neutral"])

        # Middle column (solution - light blue)
        self.add_card("Solution ideas and scope", data.get("solution", []),
                     x2, y, Inches(5.4), Inches(2.25),
                     fill_color=HC_COLORS["secondary"])
        self.add_card("User value proposition", data.get("value_prop", []),
                     x2, y + 2.4, Inches(5.4), Inches(1.25),
                     fill_color=HC_COLORS["secondary"])
        self.add_card("Adoption strategy", data.get("adoption", []),
                     x2, y + 3.75, Inches(5.4), Inches(1.6),
                     fill_color=HC_COLORS["secondary"])
        self.add_card("Success metrics", data.get("metrics", []),
                     x2, y + 5.45, Inches(5.4), Inches(1.9),
                     fill_color=HC_COLORS["secondary"])

        # Right column (business - light green)
        self.add_card("Business value", data.get("business_value", []),
                     x3, y, Inches(5.4), Inches(2.05),
                     fill_color=HC_COLORS["tertiary"])
        self.add_card("Costs & constraints", data.get("costs", []),
                     x3, y + 2.15, Inches(5.4), Inches(1.25),
                     fill_color=HC_COLORS["tertiary"])
        self.add_card("Risks & assumptions", data.get("risks", []),
                     x3, y + 3.5, Inches(5.4), Inches(1.6),
                     fill_color=HC_COLORS["tertiary"])
        self.add_card("Open questions", data.get("questions", []),
                     x3, y + 5.2, Inches(5.4), Inches(1.6),
                     fill_color=HC_COLORS["tertiary"])

        return self.save()


class BusinessModelCanvas(CanvasBuilder):
    """Business Model Canvas builder"""

    def __init__(self, title: str = "Business Model Canvas"):
        super().__init__(title)

    def build(self, data: Dict) -> str:
        """
        Build business model canvas

        Args:
            data: Dictionary with keys: key_partners, key_activities,
                  value_propositions, customer_relationships, customer_segments,
                  key_resources, channels, cost_structure, revenue_streams

        Returns:
            Path to saved PPTX
        """
        # 9-section layout (simplified for 3x3 grid)
        positions = [
            (0.3, 1.3, "key_partners", HC_COLORS["neutral"]),
            (2.3, 1.3, "key_activities", HC_COLORS["neutral"]),
            (4.3, 1.3, "value_propositions", HC_COLORS["secondary"]),
            (6.3, 1.3, "customer_relationships", HC_COLORS["neutral"]),
            (8.3, 1.3, "customer_segments", HC_COLORS["neutral"]),
            (0.3, 4.2, "key_resources", HC_COLORS["neutral"]),
            (4.3, 4.2, "channels", HC_COLORS["neutral"]),
            (0.3, 7.1, "cost_structure", HC_COLORS["tertiary"]),
            (4.3, 7.1, "revenue_streams", HC_COLORS["tertiary"]),
        ]

        labels = {
            "key_partners": "Key Partners",
            "key_activities": "Key Activities",
            "value_propositions": "Value Propositions",
            "customer_relationships": "Customer Relationships",
            "customer_segments": "Customer Segments",
            "key_resources": "Key Resources",
            "channels": "Channels",
            "cost_structure": "Cost Structure",
            "revenue_streams": "Revenue Streams",
        }

        for x, y, key, color in positions:
            self.add_card(
                labels[key],
                data.get(key, []),
                x, y,
                width=1.8,
                height=2.7,
                fill_color=color
            )

        return self.save()


class LeanStartupCanvas(CanvasBuilder):
    """Lean Startup Canvas builder"""

    def __init__(self, title: str = "Lean Startup Canvas"):
        super().__init__(title)

    def build(self, data: Dict) -> str:
        """
        Build lean startup canvas

        Args:
            data: Dictionary with keys: problem, solution, unique_value_prop,
                  unfair_advantage, customer_segments, existing_alternatives,
                  channels, revenue_streams, cost_structure, metrics

        Returns:
            Path to saved PPTX
        """
        x1, x2, x3 = 0.3, 5.9, 11.5
        y = 1.3

        # Left column
        self.add_card("Problem", data.get("problem", []),
                     x1, y, fill_color=HC_COLORS["neutral"])
        self.add_card("Customer Segments", data.get("customer_segments", []),
                     x1, y + 2.0, fill_color=HC_COLORS["neutral"])

        # Middle column
        self.add_card("Solution", data.get("solution", []),
                     x2, y, fill_color=HC_COLORS["secondary"])
        self.add_card("Unique Value Prop", data.get("unique_value_prop", []),
                     x2, y + 2.0, fill_color=HC_COLORS["secondary"])
        self.add_card("Unfair Advantage", data.get("unfair_advantage", []),
                     x2, y + 4.0, fill_color=HC_COLORS["secondary"])

        # Right column
        self.add_card("Existing Alternatives", data.get("existing_alternatives", []),
                     x3, y, fill_color=HC_COLORS["tertiary"])
        self.add_card("Channels", data.get("channels", []),
                     x3, y + 2.0, fill_color=HC_COLORS["tertiary"])
        self.add_card("Revenue Streams", data.get("revenue_streams", []),
                     x3, y + 4.0, fill_color=HC_COLORS["tertiary"])

        # Bottom
        self.add_card("Key Metrics", data.get("metrics", []),
                     x1, y + 4.25, Inches(5.4), Inches(1.5),
                     fill_color=HC_COLORS["neutral"])
        self.add_card("Cost Structure", data.get("cost_structure", []),
                     x2, y + 6.25, Inches(5.4), Inches(1.5),
                     fill_color=HC_COLORS["secondary"])

        return self.save()


# Example usage
if __name__ == "__main__":
    # HC AI Opportunity Canvas example
    ai_canvas_data = {
        "target_users": [
            "HC/PHAC analysts, scientists, inspectors",
            "Decision‑makers in public health/regulatory",
            "Canadian public via improved services"
        ],
        "problems": [
            "Manual analysis slows health threat response",
            "Fragmented data impedes insights",
            "High regulatory workload"
        ],
        "alternatives": [
            "Excel/Access databases, manual reports",
            "Ad hoc automation, not enterprise‑grade",
            "Heavy reliance on SMEs for repetition"
        ],
        "contexts": [
            "National outbreak detection – early threat spotting",
            "Drug & medical device review – faster safe approvals",
            "Climate‑health impact mapping – predict, mitigate risk",
            "Inspection targeting – focus on highest risk areas",
            "Policy intelligence scanning – rapid global insight"
        ],
        "solution": [
            "Enterprise AI platform via PATH",
            "AI‑native productivity and delivery tools",
            "AI‑enablers (MLOps, RAG, model serving)",
            "Integration of spatial data/geomatics intelligence",
            "Leverage Justice (Otto) and NRC AI assets",
            "Developer Copilot for secure, compliant code"
        ],
        "value_prop": [
            "50% faster decisions",
            "Improved accuracy and consistency",
            "Free up time for high‑value work"
        ],
        "adoption": [
            "Start with highest‑value, lowest‑risk priorities",
            "AI literacy and ethics training",
            "Transparent evaluation and AIA compliance"
        ],
        "metrics": [
            "≥90% AIA before production",
            "≥80% reuse of PATH AI services",
            "≥50% faster time‑to‑insight"
        ],
        "business_value": [
            "Stronger health readiness",
            "Faster regulatory processing",
            "Higher productivity and morale",
            "Better spatial awareness for decisions",
            "Increased public trust",
            "Accelerated outcomes via GC solutions reuse",
            "Reduced delivery timelines via Developer Copilot"
        ],
        "costs": [
            "PATH infrastructure investment",
            "Compliance with GC/TBS/privacy",
            "Non‑goal: replacing human expertise"
        ],
        "risks": [
            "Risk: Low adoption if trust lacking",
            "Risk: Data quality limits value",
            "Assumption: PATH scales with demand"
        ],
        "questions": [
            "Which highest‑value, lowest‑risk priorities first?",
            "How to sustain governance capacity long‑term?",
            "Next: formalize roadmap, align with ADM/DM priorities"
        ]
    }

    # Build Opportunity Canvas
    canvas = OpportunityCanvas("Opportunity Canvas — AI at HC/PHAC")
    path = canvas.build(ai_canvas_data)
    print(f"✓ Created: {path}")

    # Build Business Model Canvas example
    bm_data = {
        "key_partners": ["Platform partners", "Data providers"],
        "key_activities": ["AI model development", "Data integration"],
        "value_propositions": ["50% faster analysis", "Better decisions"],
        "customer_relationships": ["Support teams", "Training programs"],
        "customer_segments": ["HC analysts", "Regulators", "Public"],
        "key_resources": ["AI talent", "Data infrastructure"],
        "channels": ["Web platform", "Training"],
        "cost_structure": ["Infrastructure", "Staffing", "Licensing"],
        "revenue_streams": ["Improved services", "Operational savings"]
    }

    bm_canvas = BusinessModelCanvas("HC AI - Business Model Canvas")
    path = bm_canvas.build(bm_data)
    print(f"✓ Created: {path}")
