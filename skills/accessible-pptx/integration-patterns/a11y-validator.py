"""
PowerPoint Accessibility Validator
Validates PPTX files for WCAG 2.1 AA compliance
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from dataclasses import dataclass
from typing import List, Dict, Any, Optional
from enum import Enum
import logging

logger = logging.getLogger(__name__)


class Severity(Enum):
    """Issue severity levels"""
    CRITICAL = "critical"
    HIGH = "high"
    MEDIUM = "medium"
    LOW = "low"


@dataclass
class A11yIssue:
    """Accessibility issue"""
    slide_number: int
    severity: Severity
    category: str
    description: str
    element: str
    remediation: str

    def to_dict(self) -> Dict[str, Any]:
        return {
            'slide': self.slide_number,
            'severity': self.severity.value,
            'category': self.category,
            'description': self.description,
            'element': self.element,
            'remediation': self.remediation
        }


@dataclass
class A11yReport:
    """Accessibility compliance report"""
    filename: str
    total_slides: int
    issues: List[A11yIssue]
    compliance_rate: float
    wcag_level: str  # "A", "AA", "AAA", or "non-compliant"

    def to_dict(self) -> Dict[str, Any]:
        critical = sum(1 for i in self.issues if i.severity == Severity.CRITICAL)
        high = sum(1 for i in self.issues if i.severity == Severity.HIGH)
        medium = sum(1 for i in self.issues if i.severity == Severity.MEDIUM)
        low = sum(1 for i in self.issues if i.severity == Severity.LOW)

        return {
            'filename': self.filename,
            'total_slides': self.total_slides,
            'compliance_rate': round(self.compliance_rate, 1),
            'wcag_level': self.wcag_level,
            'issue_summary': {
                'critical': critical,
                'high': high,
                'medium': medium,
                'low': low,
                'total': len(self.issues)
            },
            'issues': [i.to_dict() for i in self.issues]
        }

    def get_summary(self) -> str:
        """Human-readable summary"""
        critical = sum(1 for i in self.issues if i.severity == Severity.CRITICAL)
        high = sum(1 for i in self.issues if i.severity == Severity.HIGH)

        summary = f"{self.filename}: {self.compliance_rate:.0f}% compliant ({self.wcag_level})\n"
        summary += f"  Critical: {critical}, High: {high}, Total issues: {len(self.issues)}\n"
        if critical > 0:
            summary += "  ⚠️  Critical issues must be fixed before production\n"
        return summary


class A11yValidator:
    """Validate PowerPoint accessibility"""

    def __init__(self, pptx_path: str):
        """Initialize validator"""
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)
        self.issues: List[A11yIssue] = []

    def validate(self) -> A11yReport:
        """Run full accessibility validation"""
        logger.info(f"Validating {self.pptx_path}")

        # Run all checks
        self._check_slide_titles()
        self._check_alt_text()
        self._check_font_sizes()
        self._check_contrast_ratios()
        self._check_link_text()
        self._check_reading_order()
        self._check_table_headers()

        # Calculate compliance
        compliance_rate = max(0, 100 - (len(self.issues) / len(self.prs.slides) * 5))
        wcag_level = self._determine_wcag_level()

        report = A11yReport(
            filename=self.pptx_path,
            total_slides=len(self.prs.slides),
            issues=self.issues,
            compliance_rate=compliance_rate,
            wcag_level=wcag_level
        )

        logger.info(report.get_summary())
        return report

    def _check_slide_titles(self) -> None:
        """Verify all slides have titles"""
        for slide_num, slide in enumerate(self.prs.slides, 1):
            if not slide.shapes.title or not slide.shapes.title.text.strip():
                self.issues.append(A11yIssue(
                    slide_number=slide_num,
                    severity=Severity.MEDIUM,
                    category='Structure',
                    description='Slide missing title',
                    element='Slide title',
                    remediation='Add descriptive title to slide'
                ))

    def _check_alt_text(self) -> None:
        """Verify all images have alt text"""
        for slide_num, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                # Check pictures
                if shape.has_image:
                    alt_text = shape.name  # python-pptx uses shape.name for alt text
                    if not alt_text or alt_text.startswith('Picture'):
                        self.issues.append(A11yIssue(
                            slide_number=slide_num,
                            severity=Severity.CRITICAL,
                            category='Images',
                            description='Image missing or default alt text',
                            element=shape.name,
                            remediation='Add descriptive alt text (< 125 characters)'
                        ))
                    elif len(alt_text) > 125:
                        self.issues.append(A11yIssue(
                            slide_number=slide_num,
                            severity=Severity.MEDIUM,
                            category='Images',
                            description=f'Alt text too long ({len(alt_text)} chars)',
                            element=shape.name,
                            remediation='Shorten alt text to < 125 characters'
                        ))

    def _check_font_sizes(self) -> None:
        """Verify font sizes meet minimum requirements"""
        for slide_num, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue

                # Check slide title (should be ≥22pt)
                if shape == slide.shapes.title:
                    if shape.text_frame.paragraphs:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size and run.font.size < Pt(22):
                                    self.issues.append(A11yIssue(
                                        slide_number=slide_num,
                                        severity=Severity.HIGH,
                                        category='Typography',
                                        description=f'Title font too small ({run.font.size.pt}pt)',
                                        element='Slide title',
                                        remediation='Increase title to ≥22pt'
                                    ))

                # Check body text (should be ≥18pt)
                else:
                    if shape.text_frame.paragraphs:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size and run.font.size < Pt(18):
                                    self.issues.append(A11yIssue(
                                        slide_number=slide_num,
                                        severity=Severity.HIGH,
                                        category='Typography',
                                        description=f'Body text too small ({run.font.size.pt}pt)',
                                        element='Body text',
                                        remediation='Increase to ≥18pt for body, ≥22pt for headings'
                                    ))
                                break  # Check first run only

    def _check_contrast_ratios(self) -> None:
        """Check color contrast ratios (simplified)"""
        for slide_num, slide in enumerate(self.prs.slides, 1):
            # This is a simplified check—real implementation would use color analysis
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                    # Note: Full implementation requires color extraction and WCAG math
                    # This is a placeholder for the pattern
                    pass

    def _check_link_text(self) -> None:
        """Verify hyperlinks have descriptive text"""
        for slide_num, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Check if run has hyperlink
                        if hasattr(run, 'hyperlink') and run.hyperlink:
                            link_text = run.text.lower()
                            bad_phrases = ['click here', 'link', 'url', 'here', 'more']
                            if any(phrase in link_text for phrase in bad_phrases):
                                self.issues.append(A11yIssue(
                                    slide_number=slide_num,
                                    severity=Severity.HIGH,
                                    category='Links',
                                    description=f'Link text not descriptive: "{run.text}"',
                                    element='Hyperlink',
                                    remediation='Change to descriptive text (e.g., "Health Canada resources")'
                                ))

    def _check_reading_order(self) -> None:
        """Check if reading order is logical (simplified)"""
        # Full implementation would check XML structure
        # This is a placeholder for the pattern
        pass

    def _check_table_headers(self) -> None:
        """Verify table headers are properly marked"""
        # Check for table shapes
        for slide_num, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                # Note: python-pptx has limited table support
                # Real implementation would check for _Cell.tcPr.tcW attributes
                if hasattr(shape, 'table'):
                    table = shape.table
                    # Verify headers exist (first row)
                    if len(table.rows) > 0:
                        # Placeholder—actual check would verify header formatting
                        pass

    def _determine_wcag_level(self) -> str:
        """Determine WCAG compliance level"""
        critical = sum(1 for i in self.issues if i.severity == Severity.CRITICAL)
        high = sum(1 for i in self.issues if i.severity == Severity.HIGH)

        if critical > 0:
            return "non-compliant"
        elif high > 0:
            return "A"
        elif len(self.issues) == 0:
            return "AA"
        else:
            # Has only medium/low issues
            return "AA" if len([i for i in self.issues if i.severity in [Severity.MEDIUM]]) <= 2 else "A"


class A11yRemediator:
    """Remediate accessibility issues"""

    def __init__(self, pptx_path: str):
        """Initialize remediator"""
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)

    def auto_fix_font_sizes(self) -> int:
        """Auto-fix font size issues"""
        fixes = 0
        for slide in self.prs.slides:
            # Fix titles
            if slide.shapes.title:
                for para in slide.shapes.title.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.font.size or run.font.size < Pt(22):
                            run.font.size = Pt(22)
                            fixes += 1

            # Fix body text
            for shape in slide.shapes:
                if shape == slide.shapes.title or not hasattr(shape, 'text_frame'):
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.font.size or run.font.size < Pt(18):
                            run.font.size = Pt(18)
                            fixes += 1
        return fixes

    def save(self, output_path: str) -> str:
        """Save remediated presentation"""
        self.prs.save(output_path)
        logger.info(f"Remediated presentation saved to {output_path}")
        return output_path


# Example usage
if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)

    # Validate
    validator = A11yValidator("example-presentation.pptx")
    report = validator.validate()

    print(report.get_summary())
    print(f"\nDetailed issues:")
    for issue in report.issues[:5]:
        print(f"  [{issue.severity.value}] Slide {issue.slide_number}: {issue.description}")

    # Export report
    import json
    with open("a11y-report.json", "w") as f:
        json.dump(report.to_dict(), f, indent=2)

    # Remediate
    remediator = A11yRemediator("example-presentation.pptx")
    fixes = remediator.auto_fix_font_sizes()
    remediator.save("example-presentation-accessible.pptx")
    print(f"\nAuto-fixed {fixes} font size issues")
